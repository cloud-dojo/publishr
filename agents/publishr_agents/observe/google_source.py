"""実Google API 観測ソース（PUBLISHR_OBSERVE=google・隔離）。

Drive/Calendar/Tasks を OAuth ユーザー資格情報で読み、ObservationBundle を組み立てる。
google-api-python-client / google-auth-oauthlib（`google` extra）が必要。資格情報は
OAuth ブートストラップ（scripts/google_oauth_bootstrap.py）が作るトークンJSONから読む。

決定的後処理（±14日窓・トリム・タスク絞り）は transform に委ね、ここは I/O とマッピング
（google_mapping）に専念する。レスポンス→モデルのマッピングは google_mapping で別途テスト。
"""

from __future__ import annotations

import logging
import os
from datetime import datetime, timedelta
from pathlib import Path

from publishr_schema import ObservationBundle, ReadingFB, User

from .google_mapping import map_calendar_event, map_drive_file, map_task_item
from .transform import DEFAULT_WINDOW_DAYS, build_observation_bundle, folder_label_map

logger = logging.getLogger(__name__)

# 3ソース・同一 OAuth・読み取り専用スコープ
ALL_SCOPES = {
    "drive": "https://www.googleapis.com/auth/drive.readonly",
    "calendar": "https://www.googleapis.com/auth/calendar.readonly",
    "tasks": "https://www.googleapis.com/auth/tasks.readonly",
}


def resolve_scopes() -> list[str]:
    """要求する OAuth スコープ。既定は3ソース全部。

    `PUBLISHR_GOOGLE_SCOPES=calendar,tasks` のように絞れる（`drive.readonly` は Google の
    restricted スコープで verification 要求が出るため、デモ検証では外せるようにする）。
    bootstrap と observe で同じ値を使うこと（同意トークンと要求スコープを一致させる）。
    """
    raw = os.environ.get("PUBLISHR_GOOGLE_SCOPES")
    if not raw:
        return list(ALL_SCOPES.values())
    keys = [k.strip() for k in raw.split(",") if k.strip()]
    unknown = [k for k in keys if k not in ALL_SCOPES]
    if unknown:
        raise SystemExit(f"未知のスコープキー: {unknown}（有効: {list(ALL_SCOPES)}）")
    return [ALL_SCOPES[k] for k in keys]


# 後方互換（既定=全ソース）。実際の要求は resolve_scopes() を使う。
SCOPES = list(ALL_SCOPES.values())

DEFAULT_TOKEN_PATH = ".secrets/google_token.json"
DEFAULT_TOKEN_DIR = ".secrets/oauth_tokens"


def token_path() -> Path:
    return Path(os.environ.get("PUBLISHR_GOOGLE_TOKEN", DEFAULT_TOKEN_PATH))


def token_dir() -> Path:
    """Web OAuth 経路の per-uid トークン保存ディレクトリ（gitignore 済 .secrets 配下）。

    BFF の callback（token_store）が書き、observe（load_credentials_for_uid）が読む共有パス。
    CLI ブートストラップの単一ファイル（token_path）とは別管理。
    """
    return Path(os.environ.get("PUBLISHR_GOOGLE_TOKEN_DIR", DEFAULT_TOKEN_DIR))


def per_uid_token_path(uid: str) -> Path:
    """uid ごとのトークンJSONパス。uid をサニタイズしてディレクトリ脱出を防ぐ。"""
    safe = "".join(c if (c.isalnum() or c in "-_") else "_" for c in uid)
    return token_dir() / f"{safe}.json"


def load_credentials_for_uid(uid: str):
    """Web OAuth で保存した per-uid トークンから資格情報を読み、期限切れなら refresh する。

    callback（token_store）が `per_uid_token_path(uid)` に書いたトークンを使う。CLI 単一ファイル
    経路（load_credentials）とは別。未配置は FileNotFoundError（パス検査が先＝google extra 非依存）。
    """
    path = per_uid_token_path(uid)
    if not path.exists():
        raise FileNotFoundError(
            f"uid={uid} の Google トークンがありません: {path}。"
            " 先に UI の『Googleアカウントで連携する』→ /api/auth/google/callback を通してください。"
        )
    from google.auth.transport.requests import Request
    from google.oauth2.credentials import Credentials

    creds = Credentials.from_authorized_user_file(str(path), resolve_scopes())
    if not creds.valid and creds.refresh_token:
        creds.refresh(Request())
        path.write_text(creds.to_json(), encoding="utf-8")
        try:
            path.chmod(0o600)  # refresh 書き戻しでも本人のみ読める権限を維持
        except OSError:
            pass
    return creds


def load_credentials():
    """保存済みトークンJSONから資格情報を読み、期限切れなら refresh する。"""
    from google.auth.transport.requests import Request
    from google.oauth2.credentials import Credentials

    path = token_path()
    if not path.exists():
        raise FileNotFoundError(
            f"Google トークンがありません: {path}。"
            " 先に `uv run python scripts/google_oauth_bootstrap.py` で OAuth 同意を済ませてください。"
        )
    creds = Credentials.from_authorized_user_file(str(path), resolve_scopes())
    if not creds.valid and creds.refresh_token:
        creds.refresh(Request())
        path.write_text(creds.to_json(), encoding="utf-8")
    return creds


# Drive 上にアップロードされた Office(OOXML) ファイル。Google ネイティブ(vnd.google-apps.*)は
# export 経路、これらは get_media で生バイトを取り _extract_office_text でローカル解析する。
DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_OFFICE_MIMES = {DOCX_MIME: "docx", XLSX_MIME: "xlsx", PPTX_MIME: "pptx"}


def _extract_office_text(mime: str, data: bytes) -> str:
    """Office(OOXML) 生バイト列 → 本文テキスト（python-docx/openpyxl/python-pptx・lazy import）。

    非対応 mime・解析失敗は空文字に縮退（観測は本文無しでも継続）。トリムは transform 側。
    解析ライブラリは `google` extra。テストはバイト列直接で offline 検証する。
    """
    kind = _OFFICE_MIMES.get(mime)
    if kind is None:
        return ""
    import io

    try:
        buf = io.BytesIO(data)
        if kind == "docx":
            import docx

            doc = docx.Document(buf)
            return "\n".join(p.text for p in doc.paragraphs if p.text)
        if kind == "xlsx":
            import openpyxl

            wb = openpyxl.load_workbook(buf, read_only=True, data_only=True)
            try:
                lines: list[str] = []
                for ws in wb.worksheets:
                    for row in ws.iter_rows(values_only=True):
                        cells = [str(c) for c in row if c is not None]
                        if cells:
                            lines.append("\t".join(cells))
                return "\n".join(lines)
            finally:
                wb.close()
        if kind == "pptx":
            from pptx import Presentation

            prs = Presentation(buf)
            lines = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text_frame.text:
                        lines.append(shape.text_frame.text)
            return "\n".join(lines)
    except Exception as exc:  # noqa: BLE001 — 解析失敗は本文無しに縮退（観測は継続）
        logger.warning("office text extract failed: mime=%s err=%s", mime, exc)
        return ""
    return ""


def _extract_text(drive, file: dict) -> str:
    """Drive ファイルからテキストを抽出。Google ネイティブは export、Office は get_media→ローカル解析。

    非ネイティブ／抽出不可は空文字（本文無しでも観測としてメタは残す）。トリムは transform 側。
    """
    mime = file.get("mimeType", "")
    try:
        if mime.startswith("application/vnd.google-apps."):
            if mime == "application/vnd.google-apps.spreadsheet":
                export_mime = "text/csv"
            else:
                export_mime = "text/plain"
            data = drive.files().export(fileId=file["id"], mimeType=export_mime).execute()
            return data.decode("utf-8", errors="replace") if isinstance(data, bytes) else str(data)
        if mime in _OFFICE_MIMES:
            data = drive.files().get_media(fileId=file["id"]).execute()
            return _extract_office_text(mime, data if isinstance(data, bytes) else bytes(data))
        if mime.startswith("text/"):
            data = drive.files().get_media(fileId=file["id"]).execute()
            return data.decode("utf-8", errors="replace") if isinstance(data, bytes) else str(data)
    except Exception as exc:  # noqa: BLE001 — 抽出失敗は本文無しに縮退（観測は継続）
        logger.warning("drive text extract failed: file=%s mime=%s err=%s", file.get("id"), mime, exc)
        return ""
    return ""


class GoogleObservationSource:
    """実Google API から観測束を生成。資格情報は未指定なら token JSON から読む。"""

    def __init__(self, credentials=None):
        self._creds = credentials

    def collect(self, user: User, *, now: datetime) -> ObservationBundle:
        from googleapiclient.discovery import build

        creds = self._creds or load_credentials()
        # 付与スコープに無いソースは取得しない（例: drive を外した同意で drive API を叩くと 403）。
        granted = set(resolve_scopes())
        use_drive = ALL_SCOPES["drive"] in granted
        use_calendar = ALL_SCOPES["calendar"] in granted
        use_tasks = ALL_SCOPES["tasks"] in granted

        cs = user.connected_sources
        drive_files = []
        if use_drive and cs and cs.drive and cs.drive.enabled:
            drive = build("drive", "v3", credentials=creds, cache_discovery=False)
            drive_files = self._fetch_drive(drive, cs)
        calendar_events = []
        if use_calendar and cs and cs.calendar and cs.calendar.enabled:
            calendar = build("calendar", "v3", credentials=creds, cache_discovery=False)
            calendar_events = self._fetch_calendar(calendar, cs, now)
        tasks_items = []
        if use_tasks and cs and cs.tasks and cs.tasks.enabled:
            tasks = build("tasks", "v1", credentials=creds, cache_discovery=False)
            tasks_items = self._fetch_tasks(tasks)

        # readingFB は STEP0 では空（§2）。集約は I-9（Firestore/BFF）で別途。
        return build_observation_bundle(
            user_id=user.id,
            now=now,
            drive_files=drive_files,
            calendar_events=calendar_events,
            tasks_items=tasks_items,
            reading_fb=ReadingFB(),
        )

    def _fetch_drive(self, drive, cs):
        label_of = folder_label_map(cs)
        out = []
        for fid in cs.drive.folder_ids:
            # folderId は Picker 由来の不透明IDのみ想定。クォート等を含むものは q を壊すので弾く。
            if "'" in fid or "\\" in fid:
                logger.warning("skipping suspicious drive folderId: %r", fid)
                continue
            resp = (
                drive.files()
                .list(
                    q=f"'{fid}' in parents and trashed = false",
                    fields="files(id,name,mimeType,modifiedTime)",
                    pageSize=50,
                )
                .execute()
            )
            for f in resp.get("files", []):
                out.append(map_drive_file(f, label_of.get(fid, ""), _extract_text(drive, f)))
        return out

    def _fetch_calendar(self, calendar, cs, now: datetime):
        time_min = (now - timedelta(days=DEFAULT_WINDOW_DAYS)).isoformat()
        time_max = (now + timedelta(days=DEFAULT_WINDOW_DAYS)).isoformat()
        calendar_ids = cs.calendar.calendar_ids or ["primary"]
        out = []
        for cid in calendar_ids:
            resp = (
                calendar.events()
                .list(
                    calendarId=cid,
                    timeMin=time_min,
                    timeMax=time_max,
                    singleEvents=True,
                    orderBy="startTime",
                    maxResults=100,
                )
                .execute()
            )
            out.extend(map_calendar_event(e) for e in resp.get("items", []))
        return out

    def _fetch_tasks(self, tasks):
        out = []
        tasklists = tasks.tasklists().list(maxResults=20).execute()
        for lst in tasklists.get("items", []):
            resp = (
                tasks.tasks()
                .list(tasklist=lst["id"], showCompleted=True, showHidden=False, maxResults=100)
                .execute()
            )
            out.extend(map_task_item(t) for t in resp.get("items", []))
        return out
